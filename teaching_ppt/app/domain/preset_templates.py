"""内置幻灯片外观预设（多主题配色）。

生成时在每张幻灯片上应用背景色与标题、正文样式（不依赖占位参考页）。
"""

from __future__ import annotations

from dataclasses import dataclass

from pptx.dml.color import RGBColor
from pptx.util import Pt


@dataclass(frozen=True)
class PresetVisual:
    cover_bg: tuple[int, int, int]
    content_bg: tuple[int, int, int]
    cover_title_rgb: tuple[int, int, int]
    cover_sub_rgb: tuple[int, int, int]
    content_title_rgb: tuple[int, int, int]
    body_rgb: tuple[int, int, int]
    cover_title_pt: int
    cover_sub_pt: int
    content_title_pt: int
    body_pt: int


_PRESETS: dict[str, PresetVisual] = {
    "simple": PresetVisual(
        cover_bg=(255, 255, 255),
        content_bg=(248, 249, 251),
        cover_title_rgb=(33, 37, 41),
        cover_sub_rgb=(107, 114, 128),
        content_title_rgb=(17, 24, 39),
        body_rgb=(55, 65, 81),
        cover_title_pt=40,
        cover_sub_pt=18,
        content_title_pt=32,
        body_pt=20,
    ),
    "cartoon": PresetVisual(
        cover_bg=(255, 252, 235),
        content_bg=(255, 250, 240),
        cover_title_rgb=(230, 90, 40),
        cover_sub_rgb=(120, 80, 40),
        content_title_rgb=(45, 110, 180),
        body_rgb=(60, 50, 45),
        cover_title_pt=44,
        cover_sub_pt=20,
        content_title_pt=34,
        body_pt=22,
    ),
    "academic": PresetVisual(
        cover_bg=(255, 255, 255),
        content_bg=(252, 252, 252),
        cover_title_rgb=(0, 51, 102),
        cover_sub_rgb=(71, 85, 105),
        content_title_rgb=(0, 51, 102),
        body_rgb=(30, 41, 59),
        cover_title_pt=38,
        cover_sub_pt=18,
        content_title_pt=30,
        body_pt=19,
    ),
    # 扩展：常见课堂 / 汇报场景
    "forest": PresetVisual(
        cover_bg=(240, 253, 244),
        content_bg=(236, 253, 245),
        cover_title_rgb=(22, 101, 52),
        cover_sub_rgb=(5, 150, 105),
        content_title_rgb=(21, 128, 61),
        body_rgb=(31, 41, 55),
        cover_title_pt=40,
        cover_sub_pt=18,
        content_title_pt=31,
        body_pt=20,
    ),
    "ocean": PresetVisual(
        cover_bg=(236, 254, 255),
        content_bg=(240, 253, 250),
        cover_title_rgb=(14, 116, 144),
        cover_sub_rgb=(13, 148, 136),
        content_title_rgb=(15, 118, 110),
        body_rgb=(51, 65, 85),
        cover_title_pt=40,
        cover_sub_pt=18,
        content_title_pt=31,
        body_pt=20,
    ),
    "sunset": PresetVisual(
        cover_bg=(255, 247, 237),
        content_bg=(255, 251, 245),
        cover_title_rgb=(194, 65, 12),
        cover_sub_rgb=(180, 83, 9),
        content_title_rgb=(154, 52, 18),
        body_rgb=(71, 55, 47),
        cover_title_pt=42,
        cover_sub_pt=19,
        content_title_pt=32,
        body_pt=20,
    ),
    "tech": PresetVisual(
        cover_bg=(15, 23, 42),
        content_bg=(248, 250, 252),
        cover_title_rgb=(248, 250, 252),
        cover_sub_rgb=(148, 163, 184),
        content_title_rgb=(30, 64, 175),
        body_rgb=(51, 65, 85),
        cover_title_pt=38,
        cover_sub_pt=17,
        content_title_pt=30,
        body_pt=19,
    ),
    "elegant": PresetVisual(
        cover_bg=(250, 245, 255),
        content_bg=(252, 250, 255),
        cover_title_rgb=(109, 40, 217),
        cover_sub_rgb=(139, 92, 246),
        content_title_rgb=(91, 33, 182),
        body_rgb=(55, 48, 74),
        cover_title_pt=40,
        cover_sub_pt=18,
        content_title_pt=31,
        body_pt=20,
    ),
    "chalk_dark": PresetVisual(
        cover_bg=(29, 78, 58),
        content_bg=(34, 88, 66),
        cover_title_rgb=(254, 243, 199),
        cover_sub_rgb=(253, 230, 138),
        content_title_rgb=(252, 211, 77),
        body_rgb=(254, 249, 195),
        cover_title_pt=40,
        cover_sub_pt=18,
        content_title_pt=31,
        body_pt=20,
    ),
}

# 与 _PRESETS 键一致，供 API 错误提示与文档使用
PRESET_LABELS: dict[str, str] = {
    "simple": "简约清爽",
    "cartoon": "卡通暖色",
    "academic": "学术藏青",
    "forest": "森林绿·自然",
    "ocean": "海洋青·清爽",
    "sunset": "暖阳橙·活力",
    "tech": "科技深蓝·封面深色",
    "elegant": "典雅紫·人文",
    "chalk_dark": "墨绿黑板·浅色字",
}

VALID_PRESET_IDS: frozenset[str] = frozenset(_PRESETS.keys())

if set(PRESET_LABELS) != set(_PRESETS):
    raise RuntimeError("PRESET_LABELS keys must match _PRESETS keys")


def get_preset_visual(preset_id: str) -> PresetVisual:
    if preset_id not in _PRESETS:
        raise ValueError(f"Unknown preset: {preset_id}")
    return _PRESETS[preset_id]


def get_preset_preview_tokens() -> dict[str, dict[str, str]]:
    """供前端预览区 CSS 使用，与内置 PPT 配色一致。"""

    def hx(rgb: tuple[int, int, int]) -> str:
        return "#{:02x}{:02x}{:02x}".format(*rgb)

    out: dict[str, dict[str, str]] = {}
    for pid, vis in _PRESETS.items():
        out[pid] = {
            "coverBg": hx(vis.cover_bg),
            "contentBg": hx(vis.content_bg),
            "coverTitle": hx(vis.cover_title_rgb),
            "coverSub": hx(vis.cover_sub_rgb),
            "contentTitle": hx(vis.content_title_rgb),
            "bodyText": hx(vis.body_rgb),
        }
    return out


def _slide_solid_bg(slide, rgb: tuple[int, int, int]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def _title_style(shape, *, rgb: tuple[int, int, int], size_pt: int, bold: bool = True) -> None:
    if shape is None or not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(*rgb)
        p.font.size = Pt(size_pt)
        p.font.bold = bold


def _body_placeholder_style(slide, *, rgb: tuple[int, int, int], size_pt: int) -> None:
    if len(slide.placeholders) < 2:
        return
    ph = slide.placeholders[1]
    if not ph.has_text_frame:
        return
    for p in ph.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(*rgb)
        p.font.size = Pt(size_pt)


def apply_preset_cover_slide(slide, vis: PresetVisual) -> None:
    _slide_solid_bg(slide, vis.cover_bg)
    _title_style(slide.shapes.title, rgb=vis.cover_title_rgb, size_pt=vis.cover_title_pt)
    if len(slide.placeholders) > 1 and slide.placeholders[1].has_text_frame:
        ph = slide.placeholders[1]
        for p in ph.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(*vis.cover_sub_rgb)
            p.font.size = Pt(vis.cover_sub_pt)


def apply_preset_content_slide(slide, vis: PresetVisual) -> None:
    _slide_solid_bg(slide, vis.content_bg)
    _title_style(slide.shapes.title, rgb=vis.content_title_rgb, size_pt=vis.content_title_pt)
    _body_placeholder_style(slide, rgb=vis.body_rgb, size_pt=vis.body_pt)
