from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from app.domain.models import JobCreateBody, SlideDeckModel
from app.domain.preset_templates import (
    apply_preset_content_slide,
    apply_preset_cover_slide,
    get_preset_visual,
)
from app.domain.slide_style import merge_preset_with_style


def _fill_content_slide_text_only(slide, deck_slide) -> None:
    slide.shapes.title.text = deck_slide.title
    if len(slide.placeholders) < 2:
        return
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    items = deck_slide.bullets if deck_slide.bullets else ["（本页以标题与讲稿为主）"]
    for i, bullet in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
    if deck_slide.notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = deck_slide.notes


def _add_side_cover_image(slide, prs: Presentation, image_path: Path) -> None:
    if not image_path.is_file():
        return
    try:
        sw = prs.slide_width
        pic_w = int(sw * 0.34)
        margin = int(Inches(0.35))
        left = sw - pic_w - margin
        top = int(Inches(0.85))
        slide.shapes.add_picture(str(image_path), left, top, width=pic_w)
    except Exception:
        pass


def _add_logo_corner(slide, prs: Presentation, image_path: Path) -> None:
    if not image_path.is_file():
        return
    try:
        w = int(Inches(0.85))
        left = prs.slide_width - w - int(Inches(0.25))
        top = prs.slide_height - int(Inches(0.55))
        slide.shapes.add_picture(str(image_path), left, top, width=w)
    except Exception:
        pass


def _build_slides(
    prs: Presentation,
    body: JobCreateBody,
    deck: SlideDeckModel,
    *,
    vis,
    cover_image_path: Path | None,
    logo_image_path: Path | None,
) -> None:
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    slide0 = prs.slides.add_slide(title_layout)
    slide0.shapes.title.text = body.topic
    if len(slide0.placeholders) > 1:
        ph = slide0.placeholders[1]
        if ph.has_text_frame:
            ph.text = f"受众：{body.audience}"
    apply_preset_cover_slide(slide0, vis)
    if cover_image_path:
        _add_side_cover_image(slide0, prs, cover_image_path)
    for s in deck.slides:
        slide = prs.slides.add_slide(content_layout)
        _fill_content_slide_text_only(slide, s)
        apply_preset_content_slide(slide, vis)
        if logo_image_path:
            _add_logo_corner(slide, prs, logo_image_path)


def build_pptx(
    body: JobCreateBody,
    deck: SlideDeckModel,
    output_path: Path,
    cover_image_path: Path | None = None,
    logo_image_path: Path | None = None,
) -> None:
    base_vis = get_preset_visual(body.template_preset)
    vis = merge_preset_with_style(base_vis, body.style)
    prs = Presentation()
    _build_slides(
        prs,
        body,
        deck,
        vis=vis,
        cover_image_path=cover_image_path,
        logo_image_path=logo_image_path,
    )
    prs.save(output_path)
